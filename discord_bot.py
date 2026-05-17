import os
import asyncio
import json
from pathlib import Path

import discord
from discord import app_commands
from discord.ext import commands

# Optional: load .env if present
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

TOKEN = os.getenv("DISCORD_BOT_TOKEN")
if not TOKEN:
    raise RuntimeError("DISCORD_BOT_TOKEN not set in environment")

# Configuration – update these if needed
# Use the dedicated project‑updates channel for all bot activity
PROJECT_UPDATES_CHANNEL_ID = int(os.getenv("PROJECT_UPDATES_CHANNEL_ID", "1502599788707975198"))
MONITORED_CHANNELS = {PROJECT_UPDATES_CHANNEL_ID}
ROLE_NAME = os.getenv("BOT_ROLE_NAME", "PM Team")
KNOWLEDGE_DIR = Path(os.getenv("KNOWLEDGE_DIR", r"C:\\Users\\savio\\Desktop\\Arcaisys\\knowledge"))


intents = discord.Intents.default()
intents.message_content = True
bot = commands.Bot(command_prefix="?", intents=intents)

# ---------- Knowledge base loading ----------

from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

# Optional LLM generation (Groq) – ensure GROQ_API_KEY is set
try:
    from groq import Groq
except ImportError:
    Groq = None

# Simple text extractors for supported file types

def extract_text_from_pdf(path: Path) -> str:
    """Extract text from a PDF safely.
    Tries `pdfplumber` first (more tolerant), falls back to `PyPDF2`.
    Any error results in an empty string and a warning – the bot never crashes.
    """
    # Try pdfplumber (handles many malformed PDFs)
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            pages_text = []
            for i, page in enumerate(pdf.pages, start=1):
                try:
                    txt = page.extract_text()
                    if txt:
                        pages_text.append(txt)
                except Exception as page_err:
                    print(f"[WARN] pdfplumber page {i} error in {path}: {page_err}")
            return "\n".join(pages_text)
    except Exception as e:
        # pdfplumber not available or failed – fallback to PyPDF2
        print(f"[INFO] pdfplumber failed for {path} ({e}), falling back to PyPDF2")
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            text = []
            for i in range(len(reader.pages)):
                try:
                    page = reader.pages[i]
                    txt = page.extract_text()
                    if txt:
                        text.append(txt)
                except Exception as page_err:
                    print(f"[WARN] PDF page {i+1} parse error in {path}: {page_err}")
            return "\n".join(text)
        except Exception as err:
            print(f"[WARN] Failed to extract PDF {path}: {err}")
            return ""




def extract_text_from_pptx(path: Path) -> str:
    """Extract text from a PowerPoint safely.
    Errors on a slide are logged and the slide is skipped.
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        txt = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt.append(shape.text)
        return "\n".join(txt)
    except Exception as e:
        print(f"[WARN] Failed to extract PPTX {path}: {e}")
        return ""



def extract_text_from_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def load_documents() -> list:
    """Load all supported documents from ``KNOWLEDGE_DIR``.
    Files that raise errors are skipped so a single corrupt PDF does not block
    the whole bot.  A ``.knowledgeignore`` file (similar to ``.gitignore``)
    can be placed in the knowledge root to exclude patterns.
    """
    docs = []
    if not KNOWLEDGE_DIR.exists():
        print(f"Knowledge directory {KNOWLEDGE_DIR} does not exist yet.")
        return docs
    # Load ignore patterns if the file exists
    ignore_path = KNOWLEDGE_DIR / ".knowledgeignore"
    ignore_patterns = []
    if ignore_path.is_file():
        ignore_patterns = [line.strip() for line in ignore_path.read_text().splitlines() if line.strip() and not line.startswith("#")]
    def is_ignored(p: Path) -> bool:
        from fnmatch import fnmatch
        rel = str(p.relative_to(KNOWLEDGE_DIR))
        return any(fnmatch(rel, pat) for pat in ignore_patterns)
    for file in KNOWLEDGE_DIR.rglob("*.*"):
        if is_ignored(file):
            continue
        suffix = file.suffix.lower()
        if suffix == ".pdf":
            txt = extract_text_from_pdf(file)
        elif suffix in {".ppt", ".pptx"}:
            txt = extract_text_from_pptx(file)
        elif suffix == ".docx":
            txt = extract_text_from_docx(file)
        elif suffix in {".txt", ".md"}:
            txt = extract_text_from_txt(file)
        else:
            continue
        if txt:
            docs.append({"path": str(file), "text": txt})
    return docs

# Split documents into chunks (max 500 characters) for better retrieval

def chunk_text(text: str, max_len: int = 500) -> list:
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_len
        chunks.append(text[start:end])
        start = end
    return chunks

model = SentenceTransformer("all-MiniLM-L6-v2")
index = None
metadata = []

async def build_vector_store():
    """Build the FAISS vector store in a background task.
    The function is now defensive – if a single document cannot be processed
    it is simply ignored.  This prevents the Discord heartbeat from being
    blocked for a long time when the knowledge folder grows.
    """
    global index, metadata
    docs = load_documents()
    if not docs:
        print("No documents loaded; vector store will be empty.")
        index = None
        metadata = []
        return
    chunks = []
    for doc in docs:
        for chunk in chunk_text(doc["text"]):
            chunks.append((doc["path"], chunk))
    texts = [c[1] for c in chunks]
    # Encode in batches to avoid a huge memory spike if there are many docs
    batch_size = 5000
    embeddings_list = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        emb = model.encode(batch, convert_to_numpy=True, show_progress_bar=False)
        embeddings_list.append(emb)
    embeddings = np.concatenate(embeddings_list, axis=0) if embeddings_list else np.empty((0, model.get_sentence_embedding_dimension()))
    dim = embeddings.shape[1] if embeddings.size else 0
    if dim == 0:
        print("No embeddings could be generated.")
        index = None
        metadata = []
        return
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    metadata = [{"path": c[0], "text": c[1]} for c in chunks]
    print(f"Built FAISS index with {len(texts)} chunks.")

# ---------- Bot behavior ----------

@bot.event
async def on_ready():
    print(f"{bot.user} has connected to Discord!")
    # Build the knowledge index in the background so the heartbeat stays healthy
    bot.loop.create_task(build_vector_store())

def role_has_access(member: discord.Member) -> bool:
    return any(role.name == ROLE_NAME for role in member.roles)

def embed_query(query: str, top_k: int = 3):
    if index is None:
        return []
    q_vec = model.encode([query], convert_to_numpy=True)
    D, I = index.search(q_vec, top_k)
    results = []
    for idx in I[0]:
        if idx < len(metadata):
            results.append(metadata[idx]["text"])
    return results

@bot.event
async def on_message(message: discord.Message):
    # Ignore messages from the bot itself
    if message.author == bot.user:
        return

    # Only respond in monitored channels
    if message.channel.id not in MONITORED_CHANNELS:
        return

    # Simple trigger: bot mentioned or message starts with '?' (command prefix already handled)
    if bot.user.mentioned_in(message) or message.content.startswith("?"):
        if not role_has_access(message.author):
            await message.reply("You don't have permission to use this bot.")
            return
        # Strip mention/prefix
        content = message.clean_content
        if bot.user.mentioned_in(message):
            content = content.replace(f"@{bot.user.name}", "").strip()
        if content.startswith("?"):
            content = content[1:].strip()

        if not content:
            await message.reply("Ask me something about the company, e.g., `?What is our product roadmap?`")
            return

        # Special commands handling
        lowered = content.lower().strip()
        if lowered == "refresh":
            await message.reply("🔄 Refreshing knowledge base…")
            await build_vector_store()
            await message.reply("✅ Knowledge base refreshed.")
            return
        if lowered == "more":
            # Pagination – send next chunk if stored in recent_queries
            key = (message.author.id, message.channel.id)
            state = recent_queries.get(key)
            if not state:
                await message.reply("No previous query to continue. Ask a new question first.")
                return
            next_idx = state["last_index"]
            chunks = state["chunks"]
            if next_idx >= len(chunks):
                await message.reply("No more results.")
                return
            # send next chunk (single or batch of 3)
            batch = chunks[next_idx: next_idx + 3]
            state["last_index"] = next_idx + len(batch)
            reply = "**More info:**\n" + "\n---\n".join(batch)
            await message.reply(reply)
            return
        if lowered.startswith("set_channel "):
            if not message.author.guild_permissions.administrator:
                await message.reply("Only administrators can change the monitored channel.")
                return
            try:
                new_id = int(lowered.split()[1])
                MONITORED_CHANNELS.clear()
                MONITORED_CHANNELS.add(new_id)
                await message.reply(f"✅ Monitored channel set to <#{new_id}>.")
            except Exception as e:
                await message.reply(f"Failed to set channel: {e}")
            return
        if lowered.startswith("set_role "):
            if not message.author.guild_permissions.administrator:
                await message.reply("Only administrators can change the required role.")
                return
            global ROLE_NAME
            ROLE_NAME = content.split(" ", 1)[1].strip()
            await message.reply(f"✅ Required role set to **{ROLE_NAME}**.")
            return

        # Normal query handling
        # Retrieve relevant chunks (default top 3)
        chunks = embed_query(content, top_k=3)
        if not chunks:
            await message.reply("I couldn't find any relevant information yet. Please add documents to the knowledge folder.")
            return
        # Store for pagination
        key = (message.author.id, message.channel.id)
        recent_queries[key] = {"query": content, "chunks": chunks, "last_index": len(chunks)}

        # Generate answer with LLM if available
        if Groq:
            try:
                prompt = f"You are a helpful assistant. Use the following extracted information to answer the user's question concisely.\n\nQuestion: {content}\n\nContext:\n" + "\n---\n".join(chunks)
                client = Groq(api_key=os.getenv("GROQ_API_KEY"))
                response = client.chat.completions.create(
                    model=os.getenv("GROQ_MODEL", "gemma2-9b-it"),
                    messages=[{"role": "system", "content": "Answer the user's question based only on the provided context."}, {"role": "user", "content": prompt}],
                    temperature=0.2,
                )
                answer = response.choices[0].message.content.strip()
                await message.reply(answer)
            except Exception as e:
                await message.reply(f"LLM generation failed ({e}). Falling back to raw chunks.\n\n**Relevant info:**\n" + "\n---\n".join(chunks))
        else:
            # Fallback to raw chunks
            reply = "**Relevant info:**\n" + "\n---\n".join(chunks)
            await message.reply(reply)
    else:
        await bot.process_commands(message)

# Global in‑memory store for pagination
recent_queries: dict = {}

# Run the bot
bot.run(TOKEN)
